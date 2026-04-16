"""
ppt_export.py
-------------
Generates stakeholder-ready PowerPoint decks from pipeline data using python-pptx.

Slide types:
  - season_readiness: Three-column Green/Yellow/Red overview
  - tech_summary:     Tech-by-tech grid with bandwidth, DOH, inventory
  - rccp:<tech>:      Supply vs Demand chart + metrics for a specific tech
  - matdi:            MATDI vs Target checkpoints table
  - site_tonnage:     Tonnage by manufacturing site chart
"""

from __future__ import annotations

import io
import math
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Colour palette matching John's existing PPT
# ---------------------------------------------------------------------------

GREEN = RGBColor(0x00, 0xB0, 0x50)
YELLOW = RGBColor(0xFF, 0xC0, 0x00)
RED = RGBColor(0xFF, 0x00, 0x00)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
TEAL = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0F, 0x17, 0x2A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
BORDER_GRAY = RGBColor(0xCB, 0xD5, 0xE1)

STATUS_COLORS = {"Green": GREEN, "Yellow": YELLOW, "Red": RED}
STATUS_BG = {
    "Green": RGBColor(0xE6, 0xF9, 0xED),
    "Yellow": RGBColor(0xFF, 0xF8, 0xE1),
    "Red": RGBColor(0xFD, 0xE8, 0xE8),
}

FONT_NAME = "Calibri"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _safe(v: Any) -> str:
    """Format a numeric value as a string, handling NaN/None."""
    if v is None or (isinstance(v, float) and math.isnan(v)):
        return "—"
    if isinstance(v, float):
        if abs(v) >= 1_000_000:
            return f"{v / 1_000_000:.1f}M"
        if abs(v) >= 1_000:
            return f"{v / 1_000:.0f}K"
        if abs(v) < 1 and v != 0:
            return f"{v:.1%}"
        return f"{v:,.0f}"
    return str(v)


def _set_cell(cell, text: str, bold: bool = False, size: int = 10,
              alignment: PP_ALIGN = PP_ALIGN.LEFT,
              font_color: RGBColor = DARK_GRAY,
              fill_color: RGBColor | None = None) -> None:
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def _add_title(slide, text: str, subtitle: str = "") -> None:
    """Add a title textbox at the top of the slide."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = BLACK

    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT_NAME
        run2.font.size = Pt(11)
        run2.font.color.rgb = MID_GRAY


def _add_footer(slide) -> None:
    """Add a bottom-right footer."""
    txBox = slide.shapes.add_textbox(
        Inches(8.5), Inches(7.0), Inches(4.5), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Magnum S&OP · MRF3 2026"
    run.font.name = FONT_NAME
    run.font.size = Pt(8)
    run.font.color.rgb = MID_GRAY
    run.font.italic = True


# =========================================================================
# Slide: Season Readiness Overview
# =========================================================================

def slide_season_readiness(prs: Presentation, bandwidth: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _add_title(slide, "SEASON READINESS 2026 MRF3")
    _add_footer(slide)

    green_techs = bandwidth[bandwidth["season_readiness"] == "Green"].sort_values("bandwidth", ascending=False)
    yellow_techs = bandwidth[bandwidth["season_readiness"] == "Yellow"].sort_values("bandwidth", ascending=False)
    red_techs = bandwidth[bandwidth["season_readiness"] == "Red"].sort_values("bandwidth", ascending=False)

    columns = [
        ("GREEN SEASON", GREEN, RGBColor(0xE6, 0xF9, 0xED), green_techs),
        ("SERVICE CHALLENGES", YELLOW, RGBColor(0xFF, 0xF8, 0xE1), yellow_techs),
        ("TOP LINE RISK", RED, RGBColor(0xFD, 0xE8, 0xE8), red_techs),
    ]

    col_width = Inches(3.8)
    left_start = Inches(0.6)
    top_header = Inches(1.2)

    for i, (label, color, bg, techs_df) in enumerate(columns):
        left = left_start + i * (col_width + Inches(0.3))

        # Circle indicator
        circle = slide.shapes.add_shape(
            1, left + col_width / 2 - Inches(0.2), top_header, Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Column heading
        txBox = slide.shapes.add_textbox(left, top_header + Inches(0.5), col_width, Inches(0.35))
        p = txBox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = FONT_NAME
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY

        # Count
        txBox2 = slide.shapes.add_textbox(left, top_header + Inches(0.85), col_width, Inches(0.25))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = f"{len(techs_df)} technologies"
        run2.font.name = FONT_NAME
        run2.font.size = Pt(9)
        run2.font.color.rgb = MID_GRAY

        # Tech list
        y = top_header + Inches(1.2)
        for _, row in techs_df.iterrows():
            shape = slide.shapes.add_shape(
                1, left + Inches(0.1), y, col_width - Inches(0.2), Inches(0.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
            shape.line.color.rgb = BORDER_GRAY
            shape.line.width = Pt(0.5)

            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"{row['main_tech']}  ({row['bandwidth']:.1%})"
            run.font.name = FONT_NAME
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = DARK_GRAY

            y += Inches(0.55)


# =========================================================================
# Slide: Tech Summary Table
# =========================================================================

def slide_tech_summary(prs: Presentation, master: pd.DataFrame, bandwidth: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Summary MRF3: 2026 Season Readiness")
    _add_footer(slide)

    bw_sorted = bandwidth.sort_values("bandwidth", ascending=False)
    techs = bw_sorted["main_tech"].tolist()

    inv_col = "projected_inv_cases" if "projected_inv_cases" in master.columns else "inv_cases"

    n_cols = len(techs) + 1  # +1 for row labels
    n_rows = 5  # header + bandwidth + peak DOH + peak inv + status
    rows_data: list[list[str]] = [
        ["Tech"] + techs,
        ["Bandwidth"] + [f"{r['bandwidth']:.1%}" for _, r in bw_sorted.iterrows()],
        ["Peak DOH"] + [f"{r['peak_doh']:.1f}" for _, r in bw_sorted.iterrows()],
        ["Peak Inventory"] + [
            _safe(master[master["main_tech"] == t][inv_col].max()) for t in techs
        ],
        ["Status"] + [r["season_readiness"] for _, r in bw_sorted.iterrows()],
    ]

    table_width = min(Inches(12.5), Inches(1.5) + Inches(1.2) * len(techs))
    col_w = int((table_width - Inches(1.5)) / max(len(techs), 1))
    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.4), Inches(1.3),
        table_width, Inches(3.5),
    ).table

    table.columns[0].width = Inches(1.5)
    for c in range(1, n_cols):
        table.columns[c].width = Emu(col_w)

    for r, row_data in enumerate(rows_data):
        for c, val in enumerate(row_data):
            is_header = r == 0
            is_label = c == 0
            bold = is_header or is_label
            size = 9 if not is_header else 10
            fill = None
            fc = DARK_GRAY

            if r == 4 and c > 0:  # status row
                status = val
                fill = STATUS_BG.get(status, None)
                fc = STATUS_COLORS.get(status, DARK_GRAY)
                bold = True

            if is_header and c > 0:
                fill = LIGHT_GRAY

            _set_cell(
                table.cell(r, c), val,
                bold=bold, size=size,
                alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT,
                font_color=fc, fill_color=fill,
            )


# =========================================================================
# Slide: RCCP for a specific technology
# =========================================================================

def slide_rccp(prs: Presentation, master: pd.DataFrame, bandwidth: pd.DataFrame, tech: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tech_bw = bandwidth[bandwidth["main_tech"] == tech]
    bw_str = f"{tech_bw['bandwidth'].iloc[0]:.1%}" if len(tech_bw) else "—"
    status = tech_bw["season_readiness"].iloc[0] if len(tech_bw) else "—"
    _add_title(slide, f"{tech} — Supply vs Demand 2026", f"Bandwidth: {bw_str} ({status})")
    _add_footer(slide)

    inv_col = "projected_inv_cases" if "projected_inv_cases" in master.columns else "inv_cases"
    df = master[master["main_tech"] == tech].copy()
    df["month_dt"] = pd.to_datetime(df["month"], format="%Y-%m", errors="coerce")
    df = df.sort_values("month_dt")

    months = df["month"].tolist()
    supply_vals = df["supply_cases"].tolist()
    demand_vals = df["demand_cases"].tolist()

    chart_data = CategoryChartData()
    chart_data.categories = months
    chart_data.add_series("Supply", supply_vals)
    chart_data.add_series("Demand", demand_vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.4),
        Inches(8.5), Inches(5.0), chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = FONT_NAME

    plot = chart.plots[0]
    plot.gap_width = 80
    series_supply = plot.series[0]
    series_demand = plot.series[1]
    series_supply.format.fill.solid()
    series_supply.format.fill.fore_color.rgb = INDIGO
    series_demand.format.fill.solid()
    series_demand.format.fill.fore_color.rgb = AMBER

    # Metrics box on the right
    peak_inv = df[inv_col].max()
    peak_doh = df["doh"].max() if "doh" in df.columns else None
    total_supply = df["supply_cases"].sum()
    total_demand = df["demand_cases"].sum()

    metrics_text = (
        f"Total Supply: {_safe(total_supply)} cases\n"
        f"Total Demand: {_safe(total_demand)} cases\n"
        f"Peak Inventory: {_safe(peak_inv)} cases\n"
    )
    if peak_doh is not None:
        metrics_text += f"Peak DOH: {peak_doh:.1f} days\n"
    metrics_text += f"Bandwidth: {bw_str}\nStatus: {status}"

    txBox = slide.shapes.add_textbox(Inches(9.3), Inches(1.8), Inches(3.5), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in metrics_text.strip().split("\n"):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        run = p.add_run()
        run.text = line
        run.font.name = FONT_NAME
        run.font.size = Pt(10)
        run.font.color.rgb = DARK_GRAY
        if line.startswith("Status:"):
            run.font.bold = True
            run.font.color.rgb = STATUS_COLORS.get(status, DARK_GRAY)


# =========================================================================
# Slide: MATDI Checkpoint Comparison
# =========================================================================

def slide_matdi(prs: Presentation, matdi_vs_target: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "MATDI vs Targets", "Apr / Aug / Dec inventory-day checkpoints")
    _add_footer(slide)

    if matdi_vs_target.empty:
        return

    techs = sorted(matdi_vs_target["main_tech"].unique())
    months = sorted(matdi_vs_target["month"].unique())

    n_rows = len(techs) + 1
    n_cols = 1 + len(months) * 3  # tech + (projected, target, status) per month

    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.4), Inches(1.3),
        Inches(12.0), Inches(0.35 * n_rows + 0.5),
    ).table

    table.columns[0].width = Inches(1.6)
    sub_w = int((Inches(12.0) - Inches(1.6)) / max(len(months) * 3, 1))

    # Header row
    _set_cell(table.cell(0, 0), "Technology", bold=True, size=9, fill_color=LIGHT_GRAY)
    ci = 1
    for m in months:
        month_label = m[5:]  # "04", "08", "12"
        month_names = {"04": "Apr", "08": "Aug", "12": "Dec"}
        ml = month_names.get(month_label, m)
        _set_cell(table.cell(0, ci), f"{ml} Proj", bold=True, size=8, alignment=PP_ALIGN.CENTER, fill_color=LIGHT_GRAY)
        _set_cell(table.cell(0, ci + 1), f"{ml} Tgt", bold=True, size=8, alignment=PP_ALIGN.CENTER, fill_color=LIGHT_GRAY)
        _set_cell(table.cell(0, ci + 2), "Status", bold=True, size=8, alignment=PP_ALIGN.CENTER, fill_color=LIGHT_GRAY)
        for j in range(3):
            table.columns[ci + j].width = Emu(sub_w)
        ci += 3

    for ri, tech in enumerate(techs):
        _set_cell(table.cell(ri + 1, 0), tech, bold=True, size=9)
        ci = 1
        for m in months:
            row_data = matdi_vs_target[
                (matdi_vs_target["main_tech"] == tech) & (matdi_vs_target["month"] == m)
            ]
            if len(row_data):
                r = row_data.iloc[0]
                _set_cell(table.cell(ri + 1, ci), f"{r['projected_matdi']:.1f}", size=9, alignment=PP_ALIGN.CENTER)
                _set_cell(table.cell(ri + 1, ci + 1), f"{r['target_matdi']:.1f}", size=9, alignment=PP_ALIGN.CENTER)
                status = r["status"]
                fill = STATUS_BG.get("Green" if status == "On Track" else "Red", None)
                fc = GREEN if status == "On Track" else RED
                _set_cell(table.cell(ri + 1, ci + 2), status, bold=True, size=8,
                          alignment=PP_ALIGN.CENTER, font_color=fc, fill_color=fill)
            else:
                for j in range(3):
                    _set_cell(table.cell(ri + 1, ci + j), "—", size=9, alignment=PP_ALIGN.CENTER)
            ci += 3


# =========================================================================
# Slide: Site Tonnage
# =========================================================================

def slide_site_tonnage(prs: Presentation, site_supply: pd.DataFrame) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Tonnage by Manufacturing Site", "Monthly litons by plant")
    _add_footer(slide)

    tonnage = site_supply.groupby(["site_name", "month"], as_index=False)["supply_litons"].sum()
    tonnage = tonnage.sort_values(["site_name", "month"])

    months = sorted(tonnage["month"].unique())
    sites = sorted(tonnage["site_name"].unique())

    chart_data = CategoryChartData()
    chart_data.categories = months
    for site in sites:
        site_data = tonnage[tonnage["site_name"] == site]
        vals = []
        for m in months:
            row = site_data[site_data["month"] == m]
            vals.append(float(row["supply_litons"].iloc[0]) if len(row) else 0)
        chart_data.add_series(site, vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.4),
        Inches(11.0), Inches(5.0), chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.name = FONT_NAME

    total = tonnage["supply_litons"].sum()
    txBox = slide.shapes.add_textbox(Inches(9.5), Inches(0.3), Inches(3.5), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Total: {_safe(total)} litons"
    run.font.name = FONT_NAME
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = DARK_GRAY


# =========================================================================
# Main export function
# =========================================================================

def generate_pptx(
    master: pd.DataFrame,
    bandwidth: pd.DataFrame,
    matdi_vs_target: pd.DataFrame,
    site_supply: pd.DataFrame,
    slides: list[str] | None = None,
) -> bytes:
    """
    Generate a PowerPoint deck and return the .pptx file as bytes.

    slides: list of slide type strings, e.g.
        ["season_readiness", "tech_summary", "rccp:48oz", "matdi", "site_tonnage"]
        If None, generates the full deck.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if slides is None:
        slides = ["season_readiness", "tech_summary", "matdi", "site_tonnage"]
        # Also add RCCP for top 3 techs by bandwidth
        top_techs = bandwidth.sort_values("bandwidth", ascending=False).head(3)["main_tech"].tolist()
        for t in top_techs:
            slides.append(f"rccp:{t}")

    for spec in slides:
        if spec == "season_readiness":
            slide_season_readiness(prs, bandwidth)
        elif spec == "tech_summary":
            slide_tech_summary(prs, master, bandwidth)
        elif spec.startswith("rccp:"):
            tech = spec.split(":", 1)[1]
            if tech in master["main_tech"].unique():
                slide_rccp(prs, master, bandwidth, tech)
        elif spec == "matdi":
            slide_matdi(prs, matdi_vs_target)
        elif spec == "site_tonnage":
            slide_site_tonnage(prs, site_supply)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
